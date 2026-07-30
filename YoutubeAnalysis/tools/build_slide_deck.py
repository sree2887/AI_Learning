"""
build_slide_deck.py
Assembles a professional 8-slide .pptx from chart images and analysis data.
Dark theme: background #1a1a2e, accent #e94560, text white.

Usage:
    python tools/build_slide_deck.py
    python tools/build_slide_deck.py --analysis .tmp/analyzed_trends.json \
        --charts-dir .tmp/charts/ \
        --output .tmp/output/youtube_trends_2026-03-13.pptx
"""

import argparse
import json
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

DEFAULT_ANALYSIS = Path(__file__).parent.parent / ".tmp" / "analyzed_trends.json"
DEFAULT_CHARTS_DIR = Path(__file__).parent.parent / ".tmp" / "charts"
DEFAULT_OUTPUT_DIR = Path(__file__).parent.parent / ".tmp" / "output"

# Colors
C_BG = RGBColor(0x1a, 0x1a, 0x2e)
C_ACCENT = RGBColor(0xe9, 0x45, 0x60)
C_CARD = RGBColor(0x16, 0x21, 0x3e)
C_TEXT = RGBColor(0xea, 0xea, 0xea)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_ACCENT = RGBColor(0x0f, 0x34, 0x60)

# Slide dimensions (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG


def add_text_box(slide, text: str, left, top, width, height,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 italic=False, wrap=True) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_TEXT


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None) -> None:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()


def add_accent_bar(slide, top=Inches(1.05)) -> None:
    """Thin colored bar below the slide title area."""
    add_filled_rect(slide, Inches(0), top, SLIDE_W, Inches(0.04), C_ACCENT)


def add_slide_title(slide, title: str, subtitle: str = "") -> None:
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.8),
                 font_size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.25),
                     font_size=11, color=C_ACCENT, italic=True)
    add_accent_bar(slide)


def add_chart_image(slide, chart_path: Path, caption: str = "") -> None:
    """Add a chart image filling most of the slide content area."""
    if not chart_path.exists():
        add_text_box(slide, f"[Chart not found: {chart_path.name}]",
                     Inches(0.4), Inches(2), Inches(12), Inches(1),
                     font_size=12, color=C_ACCENT)
        return

    top = Inches(1.15)
    bottom_margin = Inches(0.35) if caption else Inches(0.2)
    img_height = SLIDE_H - top - bottom_margin

    slide.shapes.add_picture(
        str(chart_path),
        Inches(0.3), top,
        SLIDE_W - Inches(0.6), img_height
    )

    if caption:
        add_text_box(slide, caption,
                     Inches(0.4), SLIDE_H - Inches(0.35), Inches(12.5), Inches(0.3),
                     font_size=9, color=RGBColor(0x88, 0x88, 0xaa), italic=True)


# ── Slide builders ──────────────────────────────────────────────────────────

def slide_title(prs: Presentation, report_date: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)

    # Large centered title
    add_text_box(slide, "AI & AUTOMATION",
                 Inches(1), Inches(1.2), Inches(11.3), Inches(1.2),
                 font_size=54, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, "YouTube Trends Report",
                 Inches(1), Inches(2.5), Inches(11.3), Inches(0.8),
                 font_size=32, bold=False, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Date badge
    add_filled_rect(slide, Inches(5.2), Inches(3.5), Inches(2.9), Inches(0.5), C_CARD, C_ACCENT)
    add_text_box(slide, report_date,
                 Inches(5.2), Inches(3.5), Inches(2.9), Inches(0.5),
                 font_size=14, color=C_ACCENT, align=PP_ALIGN.CENTER, bold=True)

    add_text_box(slide, "Weekly Intelligence Report  •  AI & Automation Niche",
                 Inches(1), Inches(4.3), Inches(11.3), Inches(0.4),
                 font_size=13, color=RGBColor(0x88, 0x88, 0xaa), align=PP_ALIGN.CENTER)

    # Bottom accent line
    add_filled_rect(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.08), C_ACCENT)


def slide_executive_summary(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Executive Summary", "Key insights from this week's data")

    insights = data.get("executive_summary", [])
    if not insights:
        insights = ["No insights available — check analyzed_trends.json"]

    top = Inches(1.3)
    card_h = Inches(0.82)
    gap = Inches(0.12)

    for i, insight in enumerate(insights[:5]):
        y = top + i * (card_h + gap)
        # Card background
        add_filled_rect(slide, Inches(0.3), y, Inches(12.7), card_h, C_CARD)
        # Number badge
        add_filled_rect(slide, Inches(0.3), y, Inches(0.45), card_h, C_ACCENT)
        add_text_box(slide, str(i + 1),
                     Inches(0.3), y, Inches(0.45), card_h,
                     font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Insight text
        add_text_box(slide, insight,
                     Inches(0.85), y + Inches(0.1), Inches(12.0), card_h - Inches(0.2),
                     font_size=11, color=C_TEXT, wrap=True)


def slide_chart(prs: Presentation, title: str, subtitle: str,
                chart_path: Path, caption: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, title, subtitle)
    add_chart_image(slide, chart_path, caption)


def slide_content_ideas(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Recommended Content Ideas",
                    "Based on trending topics and content gaps identified this week")

    ideas = data.get("recommended_ideas", [])
    if not ideas:
        ideas = ["No content ideas generated — check analyzed_trends.json"]

    top = Inches(1.3)
    card_h = Inches(0.78)
    gap = Inches(0.15)

    for i, idea in enumerate(ideas[:6]):
        y = top + i * (card_h + gap)
        add_filled_rect(slide, Inches(0.3), y, Inches(12.7), card_h, C_CARD, C_DARK_ACCENT)

        # Icon-like number
        add_filled_rect(slide, Inches(0.3), y, Inches(0.08), card_h,
                        C_ACCENT)

        add_text_box(slide, f"  {idea}",
                     Inches(0.5), y + Inches(0.08), Inches(12.3), card_h - Inches(0.16),
                     font_size=11, color=C_TEXT, wrap=True)


# ── Main assembler ───────────────────────────────────────────────────────────

def build_deck(analysis_path: Path, charts_dir: Path, output_path: Path) -> None:
    if not analysis_path.exists():
        print(f"Error: {analysis_path} not found. Run analyze_trends.py first.")
        sys.exit(1)

    with open(analysis_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    report_date = data.get("report_date", datetime.now().strftime("%Y-%m-%d"))

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")

    # Slide 1: Title
    slide_title(prs, report_date)
    print("  ✓ Slide 1: Title")

    # Slide 2: Executive Summary
    slide_executive_summary(prs, data)
    print("  ✓ Slide 2: Executive Summary")

    # Slide 3: Trending Keywords
    slide_chart(prs,
                "Trending Keywords",
                "Top 20 keywords ranked by frequency × average view count",
                charts_dir / "trending_keywords.png",
                "Weighted score = keyword frequency × avg views / 1K")
    print("  ✓ Slide 3: Trending Keywords")

    # Slide 4: Top Videos Table
    slide_chart(prs,
                "Top Performing Videos",
                "Ranked by engagement rate (likes + comments ÷ views)",
                charts_dir / "top_videos_table.png")
    print("  ✓ Slide 4: Top Videos")

    # Slide 5: Channel Comparison
    slide_chart(prs,
                "Channel Performance",
                "Average views per video vs engagement rate across tracked channels",
                charts_dir / "channel_comparison.png")
    print("  ✓ Slide 5: Channel Comparison")

    # Slide 6: Engagement Analysis
    slide_chart(prs,
                "Engagement Rate Analysis",
                "Views vs engagement rate — bubble size proportional to like count",
                charts_dir / "engagement_analysis.png",
                "High engagement + moderate views = strong content-audience fit")
    print("  ✓ Slide 6: Engagement Analysis")

    # Slide 7: Content Gaps
    slide_chart(prs,
                "Content Gaps & Opportunities",
                "Topics with high engagement but low saturation = your best opportunities",
                charts_dir / "content_gaps.png",
                "Score = engagement strength × (1 - topic saturation)")
    print("  ✓ Slide 7: Content Gaps")

    # Slide 8: Content Ideas
    slide_content_ideas(prs, data)
    print("  ✓ Slide 8: Content Ideas")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\nDeck saved to {output_path}")


def main():
    date_str = datetime.now().strftime("%Y-%m-%d")
    default_output = DEFAULT_OUTPUT_DIR / f"youtube_trends_{date_str}.pptx"

    parser = argparse.ArgumentParser(description="Build YouTube trends slide deck")
    parser.add_argument("--analysis", type=Path, default=DEFAULT_ANALYSIS)
    parser.add_argument("--charts-dir", type=Path, default=DEFAULT_CHARTS_DIR)
    parser.add_argument("--output", type=Path, default=default_output)
    args = parser.parse_args()

    build_deck(args.analysis, args.charts_dir, args.output)


if __name__ == "__main__":
    main()
